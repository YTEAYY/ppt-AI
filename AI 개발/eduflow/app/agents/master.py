# app/agents/master.py
"""
슬라이드 마스터 (Slide Master Agent)
====================================
역할: 분석된 콘텐츠를 바탕으로 전문적인 PPT 슬라이드를 생성합니다.

주요 기능:
1. 슬라이드 구조 설계
2. 콘텐츠 배치 및 서식 지정
3. 시각적 요소(차트, 이미지) 제안
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from langchain_openai import OpenAI
from app.services.openai_client import get_llm


def create_slide(prs: Presentation, title: str, content: str, layout_index: int = 1) -> None:
    """슬라이드 생성 유틸리티"""
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    if content and len(slide.placeholders) > 1:
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = content


def generate_ppt(outline: str, script: str, output_path: Path) -> Path:
    """
    outline:  분석된 슬라이드 구조 (제목, 핵심 내용, 시각화 아이디어)
    script : 전체 발표 대본
    
    개선 사항:
    - outline을 파싱하여 구조화된 슬라이드 생성
    - script 내용을 결론 슬라이드에 활용
    - 학교 양식에 맞는 디자인 적용
    """
    prs = Presentation()
    
    # 1) 표지 슬라이드
    create_slide(prs, "EduFlow 발표자료", "학교 맞춤형 자동 생성 자료", layout_index=0)
    
    # 2) 개요 슬라이드
    create_slide(prs, "목차", _parse_outline_summary(outline), layout_index=1)
    
    # 3) 본문 슬라이드들 파싱
    slides_data = _parse_outline(outline)
    for slide_info in slides_data:
        create_slide(
            prs, 
            slide_info["title"], 
            slide_info["content"],
            layout_index=1
        )
    
    # 4) 결론 슬라이드 - script에서 핵심 메시지 추출
    conclusion = _extract_conclusion(script)
    create_slide(prs, "결론 및 핵심 메시지", conclusion, layout_index=1)
    
    # 5) 참고자료 슬라이드
    create_slide(prs, "참고자료", "본 자료는 AI를 통해 자동 생성되었습니다.", layout_index=1)
    
    prs.save(output_path)
    return output_path


def _parse_outline_summary(outline: str) -> str:
    """outline에서 목차 요약 생성"""
    lines = outline.split("\n")
    summary = []
    for line in lines:
        if line.strip():
            # | 로 구분된 경우 첫 번째 부분을 제목으로
            if "|" in line:
                title = line.split("|")[0].strip()
                summary.append(f"• {title}")
            else:
                summary.append(f"• {line.strip()}")
    return "\n".join(summary[:5])  # 최대 5개


def _parse_outline(outline: str) -> list:
    """outline을 슬라이드 데이터 리스트로 파싱"""
    slides = []
    lines = outline.split("\n")
    
    for line in lines:
        line = line.strip()
        if not line or len(line) < 3:
            continue
            
        if "|" in line:
            parts = line.split("|")
            title = parts[0].strip()
            content = "|".join(parts[1:]).strip() if len(parts) > 1 else ""
        else:
            title = line
            content = ""
        
        if title:
            slides.append({"title": title, "content": content})
    
    return slides[:8]  # 최대 8개 슬라이드


def _extract_conclusion(script: str) -> str:
    """script에서 핵심 메시지 추출"""
    if not script:
        return "발표 내용을 정리합니다."
    
    # 첫 번째 문장이나 첫 번째 단락을 결론으로
    lines = script.split("\n")
    for line in lines:
        line = line.strip()
        if len(line) > 20:  # 충분한 길이의 텍스트
            return line[:200] + "..." if len(line) > 200 else line
    
    return "오늘의 핵심 메시지를 정리합니다."
